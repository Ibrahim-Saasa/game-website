from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "files" / "AL-RUF_Project_Presentation_Improved.pptx"
BLACK = RGBColor(5, 6, 10)
PANEL = RGBColor(18, 20, 30)
PANEL_2 = RGBColor(28, 31, 45)
ROSE = RGBColor(190, 143, 151)
BLUE = RGBColor(75, 102, 218)
SKY = RGBColor(120, 198, 235)
WHITE = RGBColor(244, 244, 238)
MUTED = RGBColor(177, 181, 194)
GREEN = RGBColor(115, 210, 167)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def shape(slide, kind, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.color.rgb = line or fill
    return s

def textbox(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04); tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02); tf.vertical_anchor = valign
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.alignment = align; p.space_after = Pt(0)
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def rich_text(slide, lines, x, y, w, h, size=16, color=WHITE, spacing=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03); tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, (label, body) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(spacing)
        r = p.add_run(); r.text = label; r.font.name = "Aptos"; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = ROSE
        if body:
            r2 = p.add_run(); r2.text = body; r2.font.name = "Aptos"; r2.font.size = Pt(size); r2.font.color.rgb = color
    return box

def add_bg(slide, index, title=None, kicker=None):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = BLACK
    for i in range(0, 14):
        line = shape(slide, MSO_SHAPE.RECTANGLE, 0.0, 0.72 + i * 0.5, 13.333, 0.006, RGBColor(19, 22, 34)); line.line.fill.background()
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, BLUE); shape(slide, MSO_SHAPE.RECTANGLE, 0.16, 0, 0.06, 7.5, ROSE)
    textbox(slide, "AL-RUF  /  PROJECT PRESENTATION", 0.6, 0.25, 5, 0.25, 9, MUTED, True)
    textbox(slide, f"{index:02d}", 12.15, 0.25, 0.55, 0.25, 10, SKY, True, align=PP_ALIGN.RIGHT)
    if kicker: textbox(slide, kicker.upper(), 0.7, 1.0, 5, 0.25, 10, SKY, True)
    if title: textbox(slide, title, 0.7, 1.27, 11.8, 0.65, 29, WHITE, True, font="Aptos Display")

def add_footer(slide):
    textbox(slide, "AL-RUF Gaming Web Platform", 0.7, 7.12, 5.5, 0.2, 8, MUTED)
    textbox(slide, "FRONT-END PROJECT  |  2025", 7.5, 7.12, 5.1, 0.2, 8, MUTED, align=PP_ALIGN.RIGHT)

def add_image(slide, path, x, y, w, h):
    if not path.exists():
        return
    if path.suffix.lower() == ".webp":
        converted = ROOT / ".dist" / f"{path.stem}.png"
        converted.parent.mkdir(exist_ok=True)
        if not converted.exists():
            Image.open(path).convert("RGB").save(converted, "PNG")
        path = converted
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))

def card(slide, x, y, w, h, label, title, body, accent=BLUE):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, PANEL, PANEL_2); shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, accent)
    textbox(slide, label.upper(), x + 0.25, y + 0.19, w - 0.45, 0.22, 8, accent, True)
    textbox(slide, title, x + 0.25, y + 0.52, w - 0.45, 0.48, 16, WHITE, True, font="Aptos Display")
    textbox(slide, body, x + 0.25, y + 1.08, w - 0.45, h - 1.22, 10, MUTED)

def add_transitions():
    effects = ["fade", "push", "wipe", "split", "fade", "push", "wipe", "fade", "split", "push", "fade"]
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for slide, effect in zip(prs.slides, effects):
        transition = slide._element.makeelement(f"{{{ns}}}transition", {"spd": "slow", "advClick": "1"})
        transition.append(slide._element.makeelement(f"{{{ns}}}{effect}", {}))
        slide._element.insert(1, transition)

# 1. Title
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BLACK
add_image(slide, ROOT / "img" / "about.webp", 7.25, 0, 6.083, 7.5)
shape(slide, MSO_SHAPE.RECTANGLE, 6.3, 0, 1.7, 7.5, BLACK); shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, BLUE); shape(slide, MSO_SHAPE.RECTANGLE, 0.18, 0, 0.07, 7.5, ROSE)
textbox(slide, "AL-RUF", 0.8, 0.9, 5.4, 0.7, 42, WHITE, True, font="Aptos Display")
textbox(slide, "GAMING WEB PLATFORM", 0.84, 1.72, 5.2, 0.3, 15, SKY, True)
textbox(slide, "A multi-page front-end experience for a fictional gaming universe.", 0.84, 2.45, 5.05, 0.8, 22, WHITE, True, font="Aptos Display")
textbox(slide, "PROJECT PRESENTATION", 0.84, 4.9, 3, 0.25, 10, ROSE, True)
textbox(slide, "Concept  |  Purpose  |  Modules  |  Technology", 0.84, 5.25, 5.4, 0.3, 13, MUTED)
textbox(slide, "Presented by: Ibrahim Saasa", 0.84, 6.55, 4.2, 0.25, 10, MUTED)

# 2. Concept
slide = prs.slides.add_slide(blank); add_bg(slide, 2, "A gaming world, presented as a platform", "Project concept")
textbox(slide, "AL-RUF is designed as the digital home of Zenithal: a fictional gaming brand that brings products, community, lore, services, and collections into one navigable experience.", 0.72, 2.15, 11.8, 0.8, 20, WHITE, True, font="Aptos Display")
card(slide, 0.72, 3.45, 3.85, 1.85, "01", "Discover", "A strong landing page introduces the universe through video, visual storytelling, and clear calls to action.", ROSE)
card(slide, 4.75, 3.45, 3.85, 1.85, "02", "Explore", "Visitors move through seven focused pages instead of one crowded screen.", BLUE)
card(slide, 8.78, 3.45, 3.85, 1.85, "03", "Engage", "Filters, FAQ interaction, forms, modals, and community features make the site feel alive.", SKY)
textbox(slide, "The result is a presentable prototype that feels like the front end of a larger gaming ecosystem.", 0.72, 5.8, 11.8, 0.45, 16, MUTED); add_footer(slide)

# 3. Purpose
slide = prs.slides.add_slide(blank); add_bg(slide, 3, "Why this project exists", "Purpose and objectives")
card(slide, 0.72, 2.1, 3.75, 3.8, "PURPOSE", "From static page to usable platform", "The project demonstrates how a fictional gaming idea can be transformed into a structured, responsive, and interactive web experience suitable for real users and future backend integration.", ROSE)
rich_text(slide, [("01  ", "Create a professional gaming-themed interface"), ("02  ", "Organize content into meaningful user journeys"), ("03  ", "Use JavaScript for practical front-end interaction"), ("04  ", "Make the experience responsive across devices"), ("05  ", "Build a foundation for future database features"), ("06  ", "Demonstrate applied HTML, CSS, JavaScript, and JSON")], 5.05, 2.18, 7.3, 3.4, 16, WHITE, 9); add_footer(slide)

# 4. Journey
slide = prs.slides.add_slide(blank); add_bg(slide, 4, "One journey, seven connected destinations", "User experience")
steps = [("01", "HOME", "Brand story and visual entry point"), ("02", "PRODUCTS", "Browse and filter offerings"), ("03", "NEXUS", "Community, events, and stats"), ("04", "VAULT", "Digital collections and rarity"), ("05", "ABOUT", "Lore, mission, and identity"), ("06", "SERVICES", "Offerings and FAQ support"), ("07", "CONTACT", "Validated message form")]
for i, (num, title, body) in enumerate(steps):
    x = 0.75 + (i % 4) * 3.1; y = 2.0 + (i // 4) * 2.0; card(slide, x, y, 2.7, 1.55, num, title, body, [ROSE, BLUE, SKY, GREEN][i % 4])
textbox(slide, "Shared navigation and footer patterns keep every page connected while page-specific styles give each destination a distinct role.", 0.75, 6.25, 11.4, 0.42, 15, MUTED); add_footer(slide)

# 5. Modules
slide = prs.slides.add_slide(blank); add_bg(slide, 5, "The core modules", "Module overview")
modules = [("HOME", "Visual introduction", "Hero video, brand story, featured metagame cards, and calls to action."), ("PRODUCTS", "Product showcase", "Mock JSON data, category filters, cards, and product detail modals."), ("NEXUS", "Community hub", "Global community, battles, events, ranked seasons, and platform statistics."), ("VAULT", "Digital collections", "Collection cards with rarity filters and reasons to collect."), ("ABOUT", "Lore and mission", "The Zenithal story, ancient order, convergence, and brand values."), ("SERVICES", "Support center", "Service offerings and an expandable FAQ accordion."), ("CONTACT", "Communication", "Contact details, social links, form validation, and success modal.")]
for i, (label, title, body) in enumerate(modules):
    x = 0.72 + (i % 4) * 3.1; y = 2.0 + (i // 4) * 2.0; card(slide, x, y, 2.7, 1.72, label, title, body, [ROSE, BLUE, SKY, GREEN][i % 4])
add_footer(slide)

# 6. Products and Vault
slide = prs.slides.add_slide(blank); add_bg(slide, 6, "Discovery is built into the interface", "Feature spotlight")
card(slide, 0.72, 2.0, 5.65, 3.5, "PRODUCTS MODULE", "Browse, filter, inspect", "The Products page renders mock product records from data/products.json. Users can separate hardware, games, and merchandise, then open a modal to inspect descriptions and specifications.", BLUE)
card(slide, 6.75, 2.0, 5.65, 3.5, "VAULT MODULE", "Collect, classify, return", "The Vault page presents the collectible side of the universe. Rarity-based filtering and collection cards give users a reason to explore and revisit the page.", ROSE)
textbox(slide, "This pairing demonstrates the difference between transactional discovery and world-building discovery.", 0.72, 6.2, 11.5, 0.42, 16, SKY, True); add_footer(slide)

# 7. Engagement
slide = prs.slides.add_slide(blank); add_bg(slide, 7, "The platform is more than a catalogue", "Engagement modules")
rich_text(slide, [("NEXUS  ", "A community-facing hub with global community, epic battles, seasonal events, ranked seasons, and headline statistics."), ("ABOUT  ", "A narrative layer that turns the fictional brand into a world with lore, mission, and values."), ("SERVICES  ", "A practical support layer covering premium support, tournament hosting, enterprise solutions, learning, security, and VIP membership."), ("CONTACT  ", "A direct communication path with validation and a success state, making the prototype feel operational.")], 0.8, 2.05, 7.8, 3.9, 17, WHITE, 13)
add_image(slide, ROOT / "img" / "swordman.webp", 9.5, 2.05, 2.6, 3.9); shape(slide, MSO_SHAPE.RECTANGLE, 9.5, 5.25, 2.6, 0.7, BLUE); textbox(slide, "COMMUNITY  /  STORY  /  SUPPORT", 9.7, 5.47, 2.2, 0.25, 9, WHITE, True, align=PP_ALIGN.CENTER); add_footer(slide)

# 8. Technology
slide = prs.slides.add_slide(blank); add_bg(slide, 8, "A front end designed for extension", "Technical implementation")
card(slide, 0.72, 2.0, 3.55, 3.55, "STRUCTURE", "HTML5", "Seven separate HTML pages share a consistent navigation system, footer, semantic sections, and accessible form controls.", ROSE)
card(slide, 4.55, 2.0, 3.55, 3.55, "VISUAL SYSTEM", "CSS3", "Common, animation, responsive, and page-specific styles keep the code organized while supporting motion, blur transitions, hover states, and mobile layouts.", BLUE)
card(slide, 8.38, 2.0, 3.55, 3.55, "BEHAVIOR", "JavaScript + JSON", "Dynamic product rendering, filters, FAQ accordion behavior, form validation, modal dialogs, and mock data create an interactive prototype.", SKY)
textbox(slide, "Current architecture: static front end  ->  future backend integration with PHP/MySQL", 0.72, 6.18, 11.6, 0.42, 16, GREEN, True); add_footer(slide)

# 9. Scope
slide = prs.slides.add_slide(blank); add_bg(slide, 9, "What is complete, and what comes next", "Current scope")
card(slide, 0.72, 2.0, 5.65, 3.8, "DELIVERED", "A complete presentable prototype", "Multi-page navigation\nResponsive layouts\nAnimated gaming visual language\nProduct data and filtering\nFAQ interaction\nContact validation and success modal", GREEN)
card(slide, 6.75, 2.0, 5.65, 3.8, "LIMITATIONS", "Front end only", "No user authentication\nNo live database\nNo online payment\nNo admin dashboard\nMock product data\nExternal assets require network access", ROSE); add_footer(slide)

# 10. Future
slide = prs.slides.add_slide(blank); add_bg(slide, 10, "A strong base for the next phase", "Future scope")
future = [("01", "Backend", "PHP / Node.js services with MySQL or another database."), ("02", "Accounts", "Registration, login, profiles, saved collections, and preferences."), ("03", "Commerce", "Cart, checkout, order tracking, and inventory management."), ("04", "Community", "Live events, leaderboards, comments, and moderation tools."), ("05", "Administration", "Dashboard for products, services, FAQs, and site analytics.")]
for i, (num, title, body) in enumerate(future):
    x = 0.72 + (i % 3) * 4.05; y = 2.0 + (i // 3) * 2.0; card(slide, x, y, 3.55, 1.55, num, title, body, [BLUE, ROSE, SKY, GREEN, BLUE][i])
textbox(slide, "The prototype already establishes the information architecture; the next phase can add persistence, identity, and real transactions without redesigning the entire experience.", 0.72, 6.2, 11.6, 0.48, 16, MUTED); add_footer(slide)

# 11. Close
slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BLACK
add_image(slide, ROOT / "img" / "contact-1.webp", 7.6, 0, 5.73, 7.5); shape(slide, MSO_SHAPE.RECTANGLE, 6.8, 0, 1.45, 7.5, BLACK); shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, ROSE); shape(slide, MSO_SHAPE.RECTANGLE, 0.18, 0, 0.07, 7.5, BLUE)
textbox(slide, "THANK YOU", 0.9, 1.15, 5.5, 0.65, 38, WHITE, True, font="Aptos Display"); textbox(slide, "AL-RUF Gaming Web Platform", 0.94, 2.05, 5.5, 0.35, 18, SKY, True); textbox(slide, "A front-end concept for a connected gaming universe.", 0.94, 3.0, 5.2, 0.7, 22, WHITE, True, font="Aptos Display"); textbox(slide, "Questions and discussion", 0.94, 5.45, 4.5, 0.35, 16, ROSE, True); textbox(slide, "Ibrahim Saasa", 0.94, 6.5, 4.5, 0.25, 10, MUTED)

add_transitions()
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides)} slides")